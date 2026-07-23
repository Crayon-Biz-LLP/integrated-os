"""Hybrid document text extractor.

Extracts verbatim text from documents using local libraries (PyMuPDF, python-docx,
openpyxl, python-pptx) with Gemini vision fallback for scanned documents and images.

Architecture:
  - PDF → PyMuPDF (algorithmic, 50ms, free, 100% verbatim)
  - DOCX → python-docx (algorithmic, 50ms)
  - XLSX → openpyxl (algorithmic, flattens cells to text)
  - PPTX → python-pptx (algorithmic, extracts all slide text)
  - Images → Gemini vision fallback (scanned docs, photos)
  - Audio → Gemini audio transcription (unchanged)
  - Fallback → SYNTHESIS_MODEL (gemini-3.6-flash) via direct API call
"""

import io
from typing import Optional

from core.lib.audit_logger import audit_log_sync


def extract_text_from_pdf(file_bytes: bytes) -> Optional[str]:
    """Extract verbatim text from a PDF using PyMuPDF.

    Returns the full text content, or None if extraction fails/empty.
    PyMuPDF preserves layout better than pypdf and supports table/image extraction.
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        extracted_parts = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text and text.strip():
                extracted_parts.append(text.strip())
        doc.close()
        if not extracted_parts:
            return None
        return "\n\n".join(extracted_parts)
    except Exception as e:
        audit_log_sync("extractor", "WARNING",
                       f"PyMuPDF extraction failed: {e}")
        return None


def extract_text_from_docx(file_bytes: bytes) -> Optional[str]:
    """Extract verbatim text from a DOCX file using python-docx."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        extracted_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                extracted_parts.append(text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    extracted_parts.append(" | ".join(row_texts))
        if not extracted_parts:
            return None
        return "\n".join(extracted_parts)
    except Exception as e:
        audit_log_sync("extractor", "WARNING",
                       f"DOCX extraction failed: {e}")
        return None


def extract_text_from_xlsx(file_bytes: bytes) -> Optional[str]:
    """Extract text from an XLSX file by flattening cell contents."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        extracted_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                cell_texts = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if cell_texts:
                    rows_text.append(" | ".join(cell_texts))
            if rows_text:
                extracted_parts.append(f"[Sheet: {sheet_name}]")
                extracted_parts.extend(rows_text)
        wb.close()
        if not extracted_parts:
            return None
        return "\n".join(extracted_parts)
    except Exception as e:
        audit_log_sync("extractor", "WARNING",
                       f"XLSX extraction failed: {e}")
        return None


def extract_text_from_pptx(file_bytes: bytes) -> Optional[str]:
    """Extract text from a PPTX file from all slides."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        extracted_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
            if slide_texts:
                extracted_parts.append(f"[Slide {slide_num}]")
                extracted_parts.extend(slide_texts)
        if not extracted_parts:
            return None
        return "\n".join(extracted_parts)
    except Exception as e:
        audit_log_sync("extractor", "WARNING",
                       f"PPTX extraction failed: {e}")
        return None


def extract_text(file_bytes: bytes, mime_type: str) -> Optional[str]:
    """Try local extraction first. Returns full text or None for fallback.

    Supports: PDF, DOCX, XLSX, PPTX, text/*, image/*
    Falls back to None for images and unsupported types → caller should use Gemini.
    """
    # Text files: decode directly
    if mime_type.startswith("text/") or mime_type in (
        "application/json", "application/xml",
    ):
        try:
            text = file_bytes.decode("utf-8", errors="replace").strip()
            return text if text else None
        except Exception:
            return None

    # Audio files: cannot extract locally, return None → Gemini audio pipeline
    if mime_type.startswith("audio/"):
        return None

    # Images: cannot extract locally, return None → Gemini vision fallback
    if mime_type.startswith("image/"):
        return None

    # PDF
    if mime_type == "application/pdf":
        return extract_text_from_pdf(file_bytes)

    # DOCX
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_docx(file_bytes)

    # XLSX
    if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return extract_text_from_xlsx(file_bytes)

    # PPTX
    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file_bytes)

    # Fall back to direct decode for unknown types
    try:
        text = file_bytes.decode("utf-8", errors="replace").strip()
        return text if text else None
    except Exception:
        return None
